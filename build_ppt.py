from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    TITLE_SLIDE = 0
    BULLET_SLIDE = 1
    
    # Function to add a styled title slide
    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
        
        # Style background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 12, 21) # --bg-dark from dashboard
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 242, 255) # Cyan
        title.text_frame.paragraphs[0].font.name = 'Arial'
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        return slide

    def add_content_slide(title_text, bullet_points):
        slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
        
        # Style background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(21, 22, 33) # --bg-panel
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 242, 255) # Cyan
        title.text_frame.paragraphs[0].font.bold = True
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.color.rgb = RGBColor(255, 255, 255) # White text
            p.font.size = Pt(22)
            p.space_after = Pt(14)
            
        return slide
        

    # Build the slides
    add_title_slide("RAM Sentinel", "Next-Gen RAM Optimization & Secure Volatile Storage")

    add_content_slide("The Problem", [
        "RAM Hogging: Browsers consume gigabytes of memory.",
        "Privacy Risks: Deleted files leave persistent forensic traces on SSDs.",
        "Cloud Reliance: Most optimizers require internet telemetry, ruining privacy."
    ])

    add_content_slide("The Solution - RAM Sentinel", [
        "Intelligent: Automates memory cleanup using Windows and Browser APIs.",
        "Secure: Military-grade ephemeral storage (RAM Disk).",
        "Private: 100% Offline. Zero Telemetry. Air-Gap Ready.",
        "Dynamic Ecosystem: Combines Sandboxing, Process management, and Memory Defragmentation."
    ])

    add_content_slide("System Architecture", [
        "Frontend: Modern Web Dashboard (Glassmorphism & Cyber UI).",
        "Backend: Python + Flask (The Central Neural Brain).",
        "Core Engines:",
        "  - Process Monitor (psutil) & Memory Compactor",
        "  - Tab Purger (Playwright Automation)",
        "  - Ghost Drive (ImDisk Virtual Storage)"
    ])

    add_content_slide("Key Features: Optimization", [
        "Memory Compactor: Deframents memory & calls Windows EmptyWorkingSet API.",
        "Neural Tab Purger: Pauses inactive browser tabs automatically.",
        "VRAM Scanner: Real-time GPU caching bottleneck monitoring.",
        "Result: Instantly frees up Gigabytes of RAM without closing your apps."
    ])

    add_content_slide("Key Features: Advanced Security", [
        "Ghost Drive (The Vault): An entire drive that exists strictly in RAM.",
        "Speed: 10x Faster than NVMe SSDs (Over 6,000 MB/s).",
        "Neural Auto-Sandbox: Zero-Trust protection. Isolates any unknown process rapidly consuming RAM.",
        "Emergency Panic: Instantly wipes memory traces and shreds vault data."
    ])

    add_content_slide("WAR ROOM (Game Mode)", [
        "Maximum Performance Override:",
        "  - Suspends Windows bloatware services (Telemtry, Spooler, Indexer).",
        "  - Demotes all background apps to IDLE CPU Scheduling.",
        "  - Elevates targeted games to HIGH_PRIORITY CPU scheduling.",
        "Automatically fully restores the system state upon deactivation."
    ])

    add_content_slide("The System Dashboard", [
        "Real-Time Visualization: CPU & RAM Usage Streamed.",
        "Active Live Modules Control Panel.",
        "Plus: System 'Dynamic Island' Desktop Widget for monitoring stats outside the browser.",
        "Offline Toggle."
    ])

    add_content_slide("Conclusion", [
        "RAM Sentinel is not just a cleaner; it's a Complete Privacy Ecosystem.",
        "Developed to prove that Performance + Privacy should NOT be a tradeoff.",
        "Fully Open Source & Configurable.",
        "Questions?"
    ])
    
    file_path = "RAM_Sentinel_Final_Presentation.pptx"
    prs.save(file_path)
    print(file_path)

if __name__ == '__main__':
    create_presentation()
